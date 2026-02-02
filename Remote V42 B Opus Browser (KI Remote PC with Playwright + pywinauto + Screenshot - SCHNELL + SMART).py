#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Remote V42 B - KI Remote PC with Playwright + pywinauto + Screenshot
Version: V42 B

Neu in V42 B:
- Screenshots automatisch in Zwischenablage (beide Buttons)
- KI antwortet knapper, weniger Wiederholung
- Direkte Dokumenterstellung via Python (KEIN LibreOffice GUI nötig!)
  → create_docx: Word-Dokumente (.docx)
  → create_xlsx: Excel-Tabellen (.xlsx) 
  → create_pptx: PowerPoint-Präsentationen (.pptx)
- Mini-DOM Extraktion (aus V41)
"""

import os, sys, json, re, base64, threading, time, subprocess, platform
from io import BytesIO
import tkinter as tk
from tkinter import ttk, Text, Scrollbar, Toplevel
import requests
from PIL import Image, ImageTk
import tempfile
from collections import Counter

# Document Libraries - werden bei Bedarf installiert
DOCX_AVAILABLE = False
XLSX_AVAILABLE = False
PPTX_AVAILABLE = False

try:
    from docx import Document
    from docx.shared import Pt, Inches
    DOCX_AVAILABLE = True
except ImportError:
    pass

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment
    XLSX_AVAILABLE = True
except ImportError:
    pass

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    PPTX_AVAILABLE = True
except ImportError:
    pass

try:
    import pyautogui
    from playwright.sync_api import sync_playwright
    PLAYWRIGHT_AVAILABLE = True
except ImportError as e:
    print(f"FEHLER: {e}")
    if 'playwright' in str(e).lower():
        PLAYWRIGHT_AVAILABLE = False
    else:
        sys.exit(1)

IS_WINDOWS = platform.system() == 'Windows'

PYWINAUTO_AVAILABLE = False
if IS_WINDOWS:
    try:
        from pywinauto import Application
        PYWINAUTO_AVAILABLE = True
    except:
        pass

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
LLM_CONFIG_DIR = os.path.join(SCRIPT_DIR, "Auswahl llm")
SYSTEM_PROMPT_FILE = os.path.join(SCRIPT_DIR, "system_prompt_gui_v42b.txt")

# Einzige Stelle für den GUI-Namen – hier anpassen
APP_TITLE = "Remote V42 B - KI Remote PC with Playwright + pywinauto + Screenshot"


def install_doc_libraries():
    """Installiert fehlende Dokumentbibliotheken"""
    global DOCX_AVAILABLE, XLSX_AVAILABLE, PPTX_AVAILABLE
    
    libs = []
    if not DOCX_AVAILABLE: libs.append('python-docx')
    if not XLSX_AVAILABLE: libs.append('openpyxl')
    if not PPTX_AVAILABLE: libs.append('python-pptx')
    
    if libs:
        print(f"📦 Installiere: {', '.join(libs)}...")
        try:
            subprocess.run([sys.executable, '-m', 'pip', 'install'] + libs + ['-q'],
                          capture_output=True, timeout=120)
            # Neu importieren
            if not DOCX_AVAILABLE:
                try:
                    from docx import Document
                    DOCX_AVAILABLE = True
                except: pass
            if not XLSX_AVAILABLE:
                try:
                    from openpyxl import Workbook
                    XLSX_AVAILABLE = True
                except: pass
            if not PPTX_AVAILABLE:
                try:
                    from pptx import Presentation
                    PPTX_AVAILABLE = True
                except: pass
        except Exception as e:
            print(f"⚠️ Installation fehlgeschlagen: {e}")


class FailureTracker:
    def __init__(self):
        self.failure_counts = Counter()
    
    def record_failure(self, action_data, error_msg):
        key = f"{action_data.get('action', '')}|{action_data.get('text', '')}"
        self.failure_counts[key] += 1
    
    def record_success(self, action_data):
        key = f"{action_data.get('action', '')}|{action_data.get('text', '')}"
        if key in self.failure_counts:
            del self.failure_counts[key]
    
    def get_failure_warning(self):
        warnings = [f"  ❌ {k.split('|')[0]} {c}x fehlgeschlagen" 
                   for k, c in self.failure_counts.items() if c >= 2]
        return "\n🚨 FEHLER:\n" + "\n".join(warnings) if warnings else ""
    
    def reset(self):
        self.failure_counts.clear()


DEFAULT_GUI_PROMPT = """Du bist ein Desktop-Automatisierungs-Assistent. 
EINE JSON-AKTION PRO ANTWORT!

═══════════════════════════════════════════════════════════════════════════════
⚡ ANTWORTE KURZ UND KNAPP!
═══════════════════════════════════════════════════════════════════════════════

❌ NICHT die Aufgabe des Nutzers wiederholen!
❌ NICHT erklären was du tun WIRST - TU ES EINFACH!
❌ KEINE langen Zusammenfassungen bei "done"!
✅ NUR die JSON-Aktion ausgeben!
✅ Bei "done": Maximal 1-2 Sätze!

═══════════════════════════════════════════════════════════════════════════════
🔒 SICHERHEITSREGELN
═══════════════════════════════════════════════════════════════════════════════

⛔ NIEMALS Dateien LÖSCHEN ohne EXPLIZITE Nutzer-Anweisung!
✅ LESEN ist IMMER erlaubt
✅ NEUE Dateien erstellen ist erlaubt

═══════════════════════════════════════════════════════════════════════════════
📄 DOKUMENTE ERSTELLEN (KEIN LibreOffice nötig!)
═══════════════════════════════════════════════════════════════════════════════

Word-Dokument (.docx):
{"action": "create_docx", "path": "C:\\Users\\...\\Documents\\dokument.docx", "title": "Titel", "content": "Der Text..."}

Excel-Tabelle (.xlsx):
{"action": "create_xlsx", "path": "C:\\Users\\...\\Documents\\tabelle.xlsx", "data": [["Spalte1", "Spalte2"], ["Wert1", "Wert2"]]}

PowerPoint (.pptx):
{"action": "create_pptx", "path": "C:\\Users\\...\\Documents\\praesentation.pptx", "slides": [{"title": "Folie 1", "content": "Text"}]}

═══════════════════════════════════════════════════════════════════════════════
🎯 MINI-DOM SYSTEM
═══════════════════════════════════════════════════════════════════════════════

Du bekommst automatisch eine Liste der klickbaren Elemente im Browser.
→ Nutze die EXAKTEN Selektoren aus der Liste!

═══════════════════════════════════════════════════════════════════════════════
🎯 KRITISCHE REGELN FÜR ÜBERSICHTEN
═══════════════════════════════════════════════════════════════════════════════

⚠️ NIEMALS "done" wenn read_file FEHLGESCHLAGEN ist!
⚠️ Bei JEDER Übersicht: run_commands → read_file → done

═══════════════════════════════════════════════════════════════════════════════
🔧 AKTIONEN
═══════════════════════════════════════════════════════════════════════════════

DOKUMENTE:
{"action": "create_docx", "path": "...", "title": "...", "content": "..."}
{"action": "create_xlsx", "path": "...", "data": [[...]]}
{"action": "create_pptx", "path": "...", "slides": [{...}]}

BROWSER:
{"action": "browser_start", "url": "https://..."}
{"action": "playwright_click", "selector": "CSS"}
{"action": "playwright_type", "selector": "CSS", "text": "..."}
{"action": "playwright_get_text", "selector": "body"}
{"action": "get_dom"}

MAUS/TASTATUR:
{"action": "mouse_click", "x": 500, "y": 300}
{"action": "key", "key": "Return"}
{"action": "pywinauto_type", "text": "...", "auto_enter": false}

SYSTEM:
{"action": "run_commands", "commands": ["start notepad"]}
{"action": "read_file", "path": "%TEMP%\\datei.txt"}

KONTROLLE:
{"action": "screenshot", "reason": "..."}
{"action": "wait"}
{"action": "done", "message": "Kurz!"}

Max 30 Schritte. EINE JSON-Aktion pro Antwort! KURZ ANTWORTEN!
"""

def load_system_prompt():
    for path in [SYSTEM_PROMPT_FILE, 
                 os.path.join(SCRIPT_DIR, "system_prompt_gui_v42.txt"),
                 os.path.join(SCRIPT_DIR, "system_prompt_gui_v41.txt"),
                 os.path.join(SCRIPT_DIR, "system_prompt_gui_v40b.txt")]:
        if os.path.exists(path):
            try:
                with open(path, 'r', encoding='utf-8') as f:
                    return f.read()
            except: pass
    return DEFAULT_GUI_PROMPT

def save_system_prompt(prompt_text):
    with open(SYSTEM_PROMPT_FILE, 'w', encoding='utf-8') as f:
        f.write(prompt_text)


# ═══════════════════════════════════════════════════════════════════════════════
# DOKUMENT-ERSTELLUNG (NEU!)
# ═══════════════════════════════════════════════════════════════════════════════

def create_docx_file(path, title=None, content=""):
    """Erstellt ein Word-Dokument direkt via Python"""
    if not DOCX_AVAILABLE:
        return False, "python-docx nicht installiert"
    
    try:
        from docx import Document
        from docx.shared import Pt
        
        doc = Document()
        
        if title:
            heading = doc.add_heading(title, level=1)
        
        # Content kann mehrzeilig sein
        for paragraph in content.split('\n\n'):
            if paragraph.strip():
                doc.add_paragraph(paragraph.strip())
        
        # Verzeichnis erstellen falls nötig
        dir_path = os.path.dirname(path)
        if dir_path and not os.path.exists(dir_path):
            os.makedirs(dir_path, exist_ok=True)
        
        doc.save(path)
        return True, f"Gespeichert: {path}"
    except Exception as e:
        return False, str(e)


def create_xlsx_file(path, data=None, sheet_name="Tabelle1"):
    """Erstellt eine Excel-Tabelle direkt via Python"""
    if not XLSX_AVAILABLE:
        return False, "openpyxl nicht installiert"
    
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font
        
        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name
        
        if data:
            for row_idx, row in enumerate(data, 1):
                for col_idx, value in enumerate(row, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    # Erste Zeile fett (Header)
                    if row_idx == 1:
                        cell.font = Font(bold=True)
        
        # Verzeichnis erstellen falls nötig
        dir_path = os.path.dirname(path)
        if dir_path and not os.path.exists(dir_path):
            os.makedirs(dir_path, exist_ok=True)
        
        wb.save(path)
        return True, f"Gespeichert: {path}"
    except Exception as e:
        return False, str(e)


def create_pptx_file(path, slides=None, title=None):
    """Erstellt eine PowerPoint-Präsentation direkt via Python"""
    if not PPTX_AVAILABLE:
        return False, "python-pptx nicht installiert"
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        prs = Presentation()
        
        if not slides:
            slides = [{"title": title or "Präsentation", "content": ""}]
        
        for slide_data in slides:
            # Titelfolie oder Inhaltsfolie
            if slides.index(slide_data) == 0:
                slide_layout = prs.slide_layouts[0]  # Titelfolie
            else:
                slide_layout = prs.slide_layouts[1]  # Titel + Inhalt
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Titel setzen
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get('title', '')
            
            # Content setzen (wenn vorhanden)
            content = slide_data.get('content', '')
            if content and len(slide.placeholders) > 1:
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx == 1:  # Body placeholder
                        placeholder.text = content
                        break
        
        # Verzeichnis erstellen falls nötig
        dir_path = os.path.dirname(path)
        if dir_path and not os.path.exists(dir_path):
            os.makedirs(dir_path, exist_ok=True)
        
        prs.save(path)
        return True, f"Gespeichert: {path}"
    except Exception as e:
        return False, str(e)


# ═══════════════════════════════════════════════════════════════════════════════
# STANDARD-FUNKTIONEN
# ═══════════════════════════════════════════════════════════════════════════════

def get_desktop_size():
    try: return pyautogui.size()
    except: return 1920, 1080

def get_mouse_position():
    try: return pyautogui.position()
    except: return 0, 0

def press_key(key):
    if '+' in key.lower():
        pyautogui.hotkey(*key.lower().split('+'))
    else:
        key_map = {'return': 'enter', 'enter': 'enter', 'tab': 'tab', 'escape': 'esc'}
        pyautogui.press(key_map.get(key.lower(), key.lower()))

def take_screenshot():
    img = pyautogui.screenshot()
    tmp_path = os.path.join(tempfile.gettempdir(), f'screenshot_v42_{int(time.time()*1000)}.png')
    img.save(tmp_path)
    buf = BytesIO()
    img.save(buf, format='PNG')
    
    # Automatisch in Zwischenablage kopieren
    copy_image_to_clipboard(img)
    
    return base64.b64encode(buf.getvalue()).decode('utf-8'), tmp_path, img


def copy_image_to_clipboard(img):
    """Kopiert PIL-Image in Windows-Zwischenablage"""
    if not IS_WINDOWS:
        return
    try:
        import win32clipboard
        output = BytesIO()
        img.convert('RGB').save(output, 'BMP')
        data = output.getvalue()[14:]  # BMP header offset
        output.close()
        
        win32clipboard.OpenClipboard()
        win32clipboard.EmptyClipboard()
        win32clipboard.SetClipboardData(win32clipboard.CF_DIB, data)
        win32clipboard.CloseClipboard()
    except ImportError:
        # win32clipboard nicht installiert - versuche mit subprocess
        try:
            tmp = os.path.join(tempfile.gettempdir(), '_clipboard_tmp.bmp')
            img.convert('RGB').save(tmp, 'BMP')
            # PowerShell-Methode
            subprocess.run([
                'powershell', '-command',
                f'Add-Type -AssemblyName System.Windows.Forms; [System.Windows.Forms.Clipboard]::SetImage([System.Drawing.Image]::FromFile("{tmp}"))'
            ], capture_output=True, creationflags=subprocess.CREATE_NO_WINDOW)
        except:
            pass
    except Exception:
        pass

def mouse_click(x=None, y=None, button='left', clicks=1):
    try:
        if x is not None and y is not None:
            pyautogui.click(x, y, clicks=clicks, button=button)
        else:
            pyautogui.click(clicks=clicks, button=button)
        return True, None
    except Exception as e:
        return False, str(e)


def load_llm_configs(config_dir):
    configs = {}
    if not os.path.exists(config_dir): return configs
    for filename in os.listdir(config_dir):
        if filename.endswith('.txt'):
            try:
                with open(os.path.join(config_dir, filename), 'r', encoding='utf-8') as f:
                    content = f.read()
                config = {}
                for mt, key in [('URL', 'url'), ('API Key', 'api_key'), ('LLM Model', 'model')]:
                    m = re.search(rf'{re.escape(mt)}:\s*(.+)', content)
                    if m: config[key] = m.group(1).strip()
                for price_label in ('Token Price', 'Kosten pro Million'):
                    m = re.search(rf'{re.escape(price_label)}:\s*([\d.,]+)', content)
                    if m:
                        try:
                            config['token_price'] = float(m.group(1).strip().replace(',', '.'))
                        except ValueError:
                            config['token_price'] = 10.0
                        break
                if 'token_price' not in config:
                    config['token_price'] = 10.0
                if 'url' in config and 'api_key' in config and 'model' in config:
                    configs[filename.replace('.txt', '')] = config
            except: pass
    return configs


def save_llm_token_price(config_dir, name, price):
    path = os.path.join(config_dir, name + '.txt')
    if not os.path.exists(path): return False
    try:
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        price_str = str(price).replace(',', '.')
        if re.search(r'Token Price:\s*[\d.,]+', content):
            content = re.sub(r'(Token Price:\s*)[\d.,]+', r'\g<1>' + price_str, content)
        elif re.search(r'Kosten pro Million:\s*[\d.,]+', content):
            content = re.sub(r'(Kosten pro Million:\s*)[\d.,]+', r'\g<1>' + price_str, content)
        else:
            content = re.sub(r'(LLM Model:\s*.+)(\r?\n)', r'\1\2Token Price: ' + price_str + r'\2', content, count=1)
        with open(path, 'w', encoding='utf-8') as f:
            f.write(content)
        return True
    except Exception:
        return False


def save_llm_config_full(config_dir, url, api_key, model, token_price):
    if not (model or '').strip():
        return False, "Modellname fehlt"
    safe_name = re.sub(r'[<>:"/\\|?*]', '_', str(model).strip())
    if not safe_name:
        return False, "Modellname ungültig"
    path = os.path.join(config_dir, safe_name + '.txt')
    try:
        content = f"""URL: {url or ''}
API Key: {api_key or ''}
LLM Model: {model.strip()}
Kosten pro Million: {token_price}
"""
        os.makedirs(config_dir, exist_ok=True)
        with open(path, 'w', encoding='utf-8') as f:
            f.write(content)
        return True, safe_name
    except Exception as e:
        return False, str(e)


def call_llm(config, messages, ss_b64=None):
    headers = {"Authorization": f"Bearer {config['api_key']}", "Content-Type": "application/json"}
    msgs = []
    for m in messages[:-1]:
        msgs.append({"role": m["role"], "content": m["content"]})
    last = messages[-1]
    if ss_b64:
        msgs.append({"role": last["role"], "content": [
            {"type": "text", "text": last["content"]},
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{ss_b64}"}}
        ]})
    else:
        msgs.append({"role": last["role"], "content": last["content"]})
    
    payload = {"model": config['model'], "messages": msgs, "max_tokens": 4000, "temperature": 0}
    try:
        r = requests.post(config['url'].rstrip('/')+'/chat/completions', headers=headers, json=payload, timeout=90)
        if r.status_code == 200:
            return r.json()['choices'][0]['message']['content'], None
        return '{"action":"done","message":"API Error"}', f"Status {r.status_code}"
    except Exception as e:
        return '{"action":"done","message":"API Error"}', str(e)


def parse_json(txt):
    if not txt: return {"action": "wait"}
    try:
        txt = re.sub(r'```(?:json)?\s*', '', txt).strip().replace('```', '')
        if '{' not in txt: return {"action": "wait"}
        
        start = txt.find('{')
        depth, in_string, escape = 0, False, False
        for i, c in enumerate(txt[start:]):
            if escape: escape = False; continue
            if c == '\\': escape = True; continue
            if c == '"': in_string = not in_string; continue
            if not in_string:
                if c == '{': depth += 1
                elif c == '}':
                    depth -= 1
                    if depth == 0:
                        try: return json.loads(txt[start:start+i+1])
                        except: break
        
        m = re.search(r'"action"\s*:\s*"([^"]+)"', txt)
        if m:
            result = {"action": m.group(1)}
            for f in ['selector', 'url', 'key', 'path', 'reason', 'message', 'text', 'title_re', 'title', 'content', 'sheet_name']:
                fm = re.search(rf'"{f}"\s*:\s*"([^"]*)"', txt)
                if fm: result[f] = fm.group(1)
            for f in ['x', 'y']:
                fm = re.search(rf'"{f}"\s*:\s*(\d+)', txt)
                if fm: result[f] = int(fm.group(1))
            for f in ['auto_enter', 'double']:
                fm = re.search(rf'"{f}"\s*:\s*(true|false)', txt, re.I)
                if fm: result[f] = fm.group(1).lower() == 'true'
            if result['action'] == 'run_commands':
                cm = re.search(r'"commands"\s*:\s*\[(.*?)\]', txt, re.DOTALL)
                if cm: result['commands'] = re.findall(r'"([^"]*)"', cm.group(1))
            # data für xlsx
            if result['action'] == 'create_xlsx':
                dm = re.search(r'"data"\s*:\s*(\[.*?\])\s*[,}]', txt, re.DOTALL)
                if dm:
                    try: result['data'] = json.loads(dm.group(1))
                    except: pass
            # slides für pptx
            if result['action'] == 'create_pptx':
                sm = re.search(r'"slides"\s*:\s*(\[.*?\])\s*[,}]', txt, re.DOTALL)
                if sm:
                    try: result['slides'] = json.loads(sm.group(1))
                    except: pass
            return result
        return {"action": "wait"}
    except:
        return {"action": "wait"}


def find_chrome_path():
    for p in [r"C:\Program Files\Google\Chrome\Application\chrome.exe",
              r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
              os.path.expandvars(r"%LOCALAPPDATA%\Google\Chrome\Application\chrome.exe")]:
        if os.path.exists(p): return p
    return "chrome"


class BrowserHelper:
    def __init__(self):
        self.playwright = self.browser = self.page = None
        self.connected = False
    
    def connect(self, port=9222):
        if not PLAYWRIGHT_AVAILABLE: return False
        for _ in range(10):
            try:
                if not self.playwright:
                    self.playwright = sync_playwright().start()
                self.browser = self.playwright.chromium.connect_over_cdp(f"http://localhost:{port}")
                ctx = self.browser.contexts
                self.page = ctx[0].pages[0] if ctx and ctx[0].pages else (ctx[0].new_page() if ctx else self.browser.new_page())
                self.connected = True
                return True
            except: time.sleep(0.5)
        return False
    
    def disconnect(self):
        try:
            if self.browser: self.browser.close()
            if self.playwright: self.playwright.stop()
        except: pass
        self.browser = self.page = self.playwright = None
        self.connected = False
    
    def get_page_info(self):
        try: return {'url': self.page.url, 'title': self.page.title()} if self.page else None
        except: return None
    
    def get_mini_dom(self, max_elements=50):
        """Extrahiert nur klickbare/interaktive Elemente"""
        if not self.page:
            return None
        
        try:
            elements = self.page.evaluate("""() => {
                const results = [];
                const selectors = 'button, a, input, select, textarea, [onclick], [role="button"], [role="link"], [type="submit"]';
                const elements = document.querySelectorAll(selectors);
                
                for (let i = 0; i < Math.min(elements.length, 100); i++) {
                    const el = elements[i];
                    if (!el.offsetParent && el.tagName !== 'INPUT') continue;
                    
                    const info = {
                        tag: el.tagName.toLowerCase(),
                        id: el.id || null,
                        classes: el.className ? el.className.split(' ').filter(c => c).slice(0, 3).join('.') : null,
                        text: (el.innerText || el.value || el.placeholder || '').trim().substring(0, 50),
                        type: el.type || null,
                        href: el.href || null,
                        name: el.name || null,
                        placeholder: el.placeholder || null
                    };
                    results.push(info);
                }
                return results;
            }""")
            
            if not elements:
                return "Keine interaktiven Elemente gefunden."
            
            lines = []
            for i, el in enumerate(elements[:max_elements], 1):
                tag = el.get('tag', '?')
                selector_parts = [tag]
                if el.get('id'):
                    selector_parts.append(f"#{el['id']}")
                elif el.get('classes'):
                    selector_parts.append(f".{el['classes']}")
                elif el.get('name'):
                    selector_parts.append(f"[name=\"{el['name']}\"]")
                
                selector = ''.join(selector_parts)
                extras = []
                if el.get('type'):
                    extras.append(f"type={el['type']}")
                if el.get('href'):
                    href = el['href'][:40] + '...' if len(el.get('href', '')) > 40 else el.get('href', '')
                    extras.append(f"href=\"{href}\"")
                if el.get('placeholder'):
                    extras.append(f"placeholder=\"{el['placeholder'][:30]}\"")
                
                text = f"\"{el['text']}\"" if el.get('text') else ""
                extra_str = f" ({', '.join(extras)})" if extras else ""
                lines.append(f"[{i}] {selector}{extra_str} {text}")
            
            return "\n".join(lines)
        
        except Exception as e:
            return f"DOM-Fehler: {str(e)[:50]}"
    
    def click(self, selector=None, text=None):
        if not self.page: return False
        try:
            if text: self.page.click(f'text="{text}"', timeout=3000)
            elif selector:
                for s in selector.split(','):
                    try: self.page.click(s.strip(), timeout=3000); return True
                    except: continue
                return False
            return True
        except: return False
    
    def type_into(self, selector, text):
        if not self.page: return False
        try:
            for s in selector.split(','):
                try: self.page.fill(s.strip(), text, timeout=3000); return True
                except: continue
            return False
        except: return False
    
    def get_text(self, selector='body'):
        try: return self.page.inner_text(selector) if self.page else None
        except: return None
    
    def scroll(self, direction='down'):
        try: self.page.evaluate(f"window.scrollBy(0, {500 if direction == 'down' else -500})"); return True
        except: return False
    
    def navigate(self, url):
        try: self.page.goto(url, wait_until='domcontentloaded', timeout=30000); return True
        except: return False


class PywinautoHelper:
    def __init__(self):
        self.app = self.window = None
        self.connected = False

    def connect(self, title=None, title_re=None):
        if not PYWINAUTO_AVAILABLE: return False
        kwargs = {"title_re": title_re} if title_re else {"title": title} if title else {}
        if not kwargs: return False
        for _ in range(3):
            try:
                self.app = Application(backend="uia").connect(**kwargs)
                self.window = self.app.top_window()
                self.connected = True
                return True
            except: time.sleep(2)
        return False

    def disconnect(self):
        self.app = self.window = None
        self.connected = False


class App:
    def __init__(self):
        self.root = tk.Tk()
        self.root.title(APP_TITLE)
        self.root.geometry("1080x900")
        self.root.configure(bg='#5a5a5a')
        
        self.run = self.stop = False
        self.cur = None
        self.llms = {}
        self.desktop_w, self.desktop_h = get_desktop_size()
        self.last_screenshot_img = None
        self.last_screenshot_b64 = None
        self.screenshot_count = self.token_count = 0
        self.token_price = 10.0
        
        self.browser = BrowserHelper()
        self.pw = PywinautoHelper()
        self.failures = FailureTracker()
        
        self.tracker_on = True
        self.tracker_window = self.tracker_label = None
        
        self.system_prompt = load_system_prompt()
        self.page_text = ""
        self.mini_dom = ""
        self._file_content = ""
        self._read_ok = None
        self.msgs = None
        
        self.gui()
        self.root.title(APP_TITLE)  # nochmal setzen, falls System/Pfad den Titel überschreibt
        self.root.bind('<F1>', lambda e: self.do_stop())
        self.root.after(100, self.load_configs)
        self.root.after(500, self.start_tracker)

    def gui(self):
        tk.Label(self.root, text=APP_TITLE, 
                font=('Arial', 13, 'bold'), bg='#5a5a5a', fg='white').pack(fill='x', pady=8)
        
        f1 = tk.Frame(self.root, bg='#5a5a5a')
        f1.pack(fill='x', padx=10, pady=5)
        tk.Label(f1, text="LLM:", width=10, bg='#5a5a5a', fg='white').pack(side='left')
        self.var = tk.StringVar(value="--")
        self.combo = ttk.Combobox(f1, textvariable=self.var, values=["--"], state='readonly', width=60)
        self.combo.pack(side='left', fill='x', expand=True, padx=5)
        self.combo.bind('<<ComboboxSelected>>', self.on_select)
        
        for label, attr, show in [("URL:", 'url_ent', ''), ("Model:", 'model_ent', ''), ("API Key:", 'key_ent', '*')]:
            f = tk.Frame(self.root, bg='#5a5a5a')
            f.pack(fill='x', padx=10, pady=2)
            tk.Label(f, text=label, width=10, bg='#5a5a5a', fg='white').pack(side='left')
            ent = tk.Entry(f, show=show, bg='#d0d0d0', fg='#222')
            ent.pack(side='left', fill='x', expand=True, padx=5)
            setattr(self, attr, ent)
        
        f_info = tk.Frame(self.root, bg='#d4a0a0')
        f_info.pack(fill='x', padx=10, pady=5)
        
        # Doc-Status anzeigen
        doc_status = []
        if DOCX_AVAILABLE: doc_status.append("📄docx")
        if XLSX_AVAILABLE: doc_status.append("📊xlsx")
        if PPTX_AVAILABLE: doc_status.append("📽pptx")
        doc_str = " ".join(doc_status) if doc_status else "⚠️Docs fehlen"
        
        self.info_lbl = tk.Label(f_info, text=f"{self.desktop_w}x{self.desktop_h} | V42 B | {doc_str}",
                                 bg='#d4a0a0', fg='#4a2828')
        self.info_lbl.pack(side='left', padx=10, pady=5)
        tk.Button(f_info, text="Prompt", command=self.edit_prompt, bg='#c70', fg='white', width=7).pack(side='left', padx=2)
        tk.Button(f_info, text="Letzter Screenshot KI", command=self.show_last_screenshot, bg='#805', fg='white', width=16).pack(side='left', padx=2)
        tk.Button(f_info, text="Screenshot", command=self.test_screenshot, bg='#5a5', fg='white', width=9).pack(side='left', padx=2)
        self.tracker_btn = tk.Button(f_info, text="Tracker", command=self.toggle_tracker, bg='#0a5', fg='white', width=7)
        self.tracker_btn.pack(side='left', padx=2)
        tk.Button(f_info, text="LLM Save", command=self.llm_save, bg='#007bff', fg='white', width=7).pack(side='left', padx=2)
        
        cf = tk.Frame(self.root, bg='#5a5a5a')
        cf.pack(fill='both', expand=True, padx=10, pady=5)
        sb = Scrollbar(cf)
        sb.pack(side='right', fill='y')
        self.txt = Text(cf, wrap=tk.WORD, yscrollcommand=sb.set, font=('Monospace', 9), bg='#1a1a1a', fg='#0f0', insertbackground='#aaa')
        self.txt.pack(side='left', fill='both', expand=True)
        sb.config(command=self.txt.yview)
        
        ef = tk.Frame(self.root, bg='#5a5a5a')
        ef.pack(fill='x', padx=10, pady=5)
        self.ent = Text(ef, font=('Arial', 11), height=3, bg='#d0d0d0', fg='#222', insertbackground='#555')
        self.ent.pack(fill='x', expand=True)
        self.ent.bind('<Return>', self.send)
        self.ent.config(state='disabled')
        
        bf = tk.Frame(self.root, bg='#5a5a5a')
        bf.pack(fill='x', padx=10, pady=5)
        self.send_btn = tk.Button(bf, text="Senden", command=self.send, state='disabled', bg='#0a5', fg='white')
        self.send_btn.pack(side='left', padx=3)
        self.stop_btn = tk.Button(bf, text="Stop", command=self.do_stop, state='disabled', bg='#a33', fg='white')
        self.stop_btn.pack(side='left', padx=3)
        tk.Button(bf, text="Neu", command=self.clear, bg='#f80', fg='white').pack(side='left', padx=3)
        tk.Button(bf, text="Copy", command=self.copy, bg='#55a', fg='white').pack(side='left', padx=3)
        tk.Label(bf, text="€/1M:", bg='#5a5a5a', fg='white').pack(side='left', padx=(10,2))
        self.price_ent = tk.Entry(bf, width=6, bg='#d0d0d0', fg='#222')
        self.price_ent.insert(0, "10.0")
        self.price_ent.pack(side='left', padx=2)
        self.price_ent.bind('<FocusOut>', self.on_price_changed)
        self.price_ent.bind('<Return>', self.on_price_changed)
        tk.Label(bf, text="Tokens:", bg='#5a5a5a', fg='white').pack(side='left', padx=(8,2))
        self.token_lbl = tk.Label(bf, text="0", bg='#5a5a5a', fg='#0ff')
        self.token_lbl.pack(side='left')
        tk.Label(bf, text="≈", bg='#5a5a5a', fg='white').pack(side='left')
        self.price_lbl = tk.Label(bf, text="0.000 €", bg='#5a5a5a', fg='#0ff')
        self.price_lbl.pack(side='left')

    def show_last_screenshot(self):
        if self.last_screenshot_img is None:
            self.log("⚠️ Noch kein KI-Screenshot vorhanden.\n")
            return
        try:
            # In Zwischenablage kopieren
            copy_image_to_clipboard(self.last_screenshot_img)
            self.log("📋 Screenshot in Zwischenablage kopiert!\n")
            
            win = Toplevel(self.root)
            win.title(f"Letzter KI-Screenshot (#{self.screenshot_count})")
            img = self.last_screenshot_img
            w, h = img.size
            r = min(1200/w, 800/h, 1)
            img_d = img.resize((int(w*r), int(h*r)), Image.Resampling.LANCZOS) if r < 1 else img
            photo = ImageTk.PhotoImage(img_d)
            lbl = tk.Label(win, image=photo)
            lbl.image = photo
            lbl.pack(padx=10, pady=10)
        except Exception as e:
            self.log(f"❌ {e}\n")

    def test_screenshot(self):
        try:
            _, _, img = take_screenshot()
            self.log("📋 Screenshot in Zwischenablage kopiert!\n")
            
            win = Toplevel(self.root)
            win.title("Aktueller Screenshot (in Zwischenablage)")
            w, h = img.size
            r = min(1200/w, 800/h, 1)
            img_d = img.resize((int(w*r), int(h*r)), Image.Resampling.LANCZOS) if r < 1 else img
            photo = ImageTk.PhotoImage(img_d)
            lbl = tk.Label(win, image=photo)
            lbl.image = photo
            lbl.pack(padx=10, pady=10)
        except Exception as e:
            self.log(f"❌ {e}\n")

    def start_tracker(self):
        self.create_tracker()
        self.update_tracker()
    
    def toggle_tracker(self):
        self.tracker_on = not self.tracker_on
        self.tracker_btn.config(text="Tracker" if self.tracker_on else "Tracker off", bg='#0a5' if self.tracker_on else '#666')
        if self.tracker_on:
            self.create_tracker()
            self.update_tracker()
        elif self.tracker_window:
            self.tracker_window.destroy()
            self.tracker_window = None

    def llm_save(self):
        url = self.url_ent.get().strip()
        model = self.model_ent.get().strip()
        api_key = self.key_ent.get().strip()
        try:
            price = float(str(self.price_ent.get()).replace(',', '.'))
        except ValueError:
            price = 10.0
        if not model:
            self.log("❌ Modellname fehlt für LLM Save.\n")
            return
        ok, msg = save_llm_config_full(LLM_CONFIG_DIR, url, api_key, model, price)
        if ok:
            self.log(f"💾 LLM \"{msg}\" gespeichert.\n")
            self.llms = load_llm_configs(LLM_CONFIG_DIR)
            if self.llms:
                self.combo['values'] = ["--"] + sorted(self.llms.keys())
            self.var.set(msg)
            if msg in self.llms:
                self.cur = self.llms[msg]
                self.ent.config(state='normal')
                self.send_btn.config(state='normal')
        else:
            self.log(f"❌ LLM Save Fehler: {msg}\n")
    
    def create_tracker(self):
        if self.tracker_window: return
        self.tracker_window = Toplevel(self.root)
        self.tracker_window.overrideredirect(True)
        self.tracker_window.attributes('-topmost', True)
        self.tracker_label = tk.Label(self.tracker_window, text="X:0 Y:0", font=('Arial', 14, 'bold'), bg='black', fg='red')
        self.tracker_label.pack(padx=5, pady=3)
    
    def update_tracker(self):
        if not self.tracker_on or not self.tracker_window: return
        try:
            x, y = get_mouse_position()
            self.tracker_window.geometry(f"+{x+15}+{y+25}")
            self.tracker_label.config(text=f"X:{x} Y:{y}")
        except: pass
        self.root.after(50, self.update_tracker)

    def edit_prompt(self):
        ed = Toplevel(self.root)
        ed.title("System-Prompt")
        ed.geometry("1000x800")
        ed.configure(bg='#2a2a3a')
        f = tk.Frame(ed)
        f.pack(fill='both', expand=True, padx=10, pady=5)
        sb = Scrollbar(f)
        sb.pack(side='right', fill='y')
        t = Text(f, wrap=tk.WORD, font=('Monospace', 10), bg='#1a1a1a', fg='#0f0', insertbackground='#aaa', yscrollcommand=sb.set)
        t.pack(side='left', fill='both', expand=True)
        sb.config(command=t.yview)
        t.insert('1.0', self.system_prompt)
        bf = tk.Frame(ed, bg='#2a2a3a')
        bf.pack(fill='x', padx=10, pady=10)
        def save():
            self.system_prompt = t.get('1.0', 'end-1c')
            save_system_prompt(self.system_prompt)
            self.log("✅ Gespeichert\n")
            ed.destroy()
        tk.Button(bf, text="Speichern", command=save, bg='#0a5', fg='white').pack(side='left', padx=5)
        tk.Button(bf, text="Standard", command=lambda: (t.delete('1.0', tk.END), t.insert('1.0', DEFAULT_GUI_PROMPT)), bg='#c70', fg='white').pack(side='left', padx=5)

    def log(self, t):
        self.txt.insert(tk.END, t)
        self.txt.see(tk.END)

    def clear(self):
        self.txt.delete('1.0', tk.END)
        self.msgs = None
        self.page_text = self._file_content = ""
        self.mini_dom = ""
        self._read_ok = None
        self.failures.reset()
        self.screenshot_count = self.token_count = 0
        self.update_tokens()
        self.log("🆕 Neuer Kontext\n")

    def update_tokens(self):
        self.token_lbl.config(text=f"{self.token_count:,}".replace(",", "."))
        self.price_lbl.config(text=f"{(self.token_count / 1_000_000) * self.token_price:.3f} €")

    def copy(self):
        self.root.clipboard_clear()
        self.root.clipboard_append(self.txt.get('1.0', tk.END))

    def load_configs(self):
        self.log("🔍 LLM Configs...\n")
        self.llms = load_llm_configs(LLM_CONFIG_DIR)
        if self.llms:
            self.combo['values'] = ["--"] + sorted(self.llms.keys())
            self.log(f"✅ {len(self.llms)} gefunden\n")
        
        # Doc-Libraries Status
        self.log(f"📄 Dokumente: docx={'✅' if DOCX_AVAILABLE else '❌'} xlsx={'✅' if XLSX_AVAILABLE else '❌'} pptx={'✅' if PPTX_AVAILABLE else '❌'}\n")

    def on_price_changed(self, e=None):
        try:
            val = self.price_ent.get().strip().replace(',', '.')
            price = float(val) if val else 10.0
            if price < 0:
                price = 10.0
            self.token_price = price
            self.price_ent.delete(0, tk.END)
            self.price_ent.insert(0, str(price))
            if self.cur is not None:
                self.cur['token_price'] = price
                n = self.var.get()
                if n != "--" and n in self.llms:
                    if save_llm_token_price(LLM_CONFIG_DIR, n, price):
                        self.log("💾 Preis gespeichert\n")
            self.update_tokens()
        except ValueError:
            self.price_ent.delete(0, tk.END)
            self.price_ent.insert(0, str(self.token_price))

    def on_select(self, e):
        n = self.var.get()
        if n != "--" and n in self.llms:
            self.cur = self.llms[n]
            self.url_ent.delete(0, tk.END); self.url_ent.insert(0, self.cur.get('url', ''))
            self.model_ent.delete(0, tk.END); self.model_ent.insert(0, self.cur.get('model', ''))
            self.key_ent.delete(0, tk.END); self.key_ent.insert(0, self.cur.get('api_key', ''))
            self.token_price = self.cur.get('token_price', 10.0)
            self.price_ent.delete(0, tk.END)
            self.price_ent.insert(0, str(self.token_price))
            self.update_tokens()
            self.log(f"\n🤖 {self.cur.get('model')} (€{self.token_price}/1M)\n✅ Bereit!\n")
            self.ent.config(state='normal')
            self.send_btn.config(state='normal')

    def send(self, e=None):
        if self.run or not self.cur: return
        cmd = self.ent.get("1.0", "end-1c").strip()
        if not cmd: return
        self.ent.delete("1.0", tk.END)
        
        force_new = cmd.upper().startswith("NEU ") or cmd.upper() == "NEU"
        if force_new:
            cmd = cmd[3:].strip() if len(cmd) > 3 else ""
            if not cmd:
                self.log("⚠️ Nach 'NEU' eine Aufgabe eingeben.\n")
                return
            self.msgs = None
            self.page_text = self._file_content = ""
            self.mini_dom = ""
            self._read_ok = None
            self.failures.reset()
            self.log("🔄 Kontext gelöscht.\n")
        
        self.log(f"\n{'='*60}\n👤 {cmd}\n{'='*60}\n")
        self.run = True
        self.stop = False
        self.send_btn.config(state='disabled')
        self.stop_btn.config(state='normal')
        self.ent.config(state='disabled')
        threading.Thread(target=self.work, args=(cmd, force_new), daemon=True).start()

    def do_stop(self):
        self.stop = True

    def work(self, cmd, force_new):
        if self.msgs is None or force_new:
            msgs = [{"role": "system", "content": self.system_prompt},
                    {"role": "user", "content": f"Aufgabe: {cmd}\n\nDENKE → HANDLE → KONTROLLIERE. EINE JSON-Aktion!"}]
            self.page_text = self._file_content = ""
            self.mini_dom = ""
            self._read_ok = None
            self.failures.reset()
            self.root.after(0, lambda: self.log("🔍 Neuer Kontext\n"))
        else:
            msgs = list(self.msgs)
            msgs.append({"role": "user", "content": f"Nachfrage: {cmd}"})
            self.root.after(0, lambda: self.log("🔍 Kontext beibehalten\n"))
        
        self.pw.disconnect()
        if PYWINAUTO_AVAILABLE: self.root.after(0, lambda: self.log("✅ pywinauto\n"))
        if PLAYWRIGHT_AVAILABLE: self.root.after(0, lambda: self.log("✅ Playwright\n"))
        
        ss_b64 = None
        
        for i in range(30):
            if self.stop:
                self.root.after(0, lambda: self.log("⏹ Stop\n"))
                break
            
            self.root.after(0, lambda i=i: self.log(f"\n[{i+1}] 🧠 Denke...\n"))
            
            last = msgs[-1]
            enhanced = last["content"]
            
            warn = self.failures.get_failure_warning()
            if warn: enhanced = warn + "\n" + enhanced
            
            mx, my = get_mouse_position()
            enhanced += f"\n\n🖱️ Maus: {mx},{my} | Bildschirm: {self.desktop_w}x{self.desktop_h}"
            
            # Dokumenten-Info
            doc_info = []
            if DOCX_AVAILABLE: doc_info.append("docx")
            if XLSX_AVAILABLE: doc_info.append("xlsx")
            if PPTX_AVAILABLE: doc_info.append("pptx")
            if doc_info:
                enhanced += f"\n📄 Verfügbare Dokument-Actions: create_{', create_'.join(doc_info)}"
            
            if self.browser.connected:
                info = self.browser.get_page_info()
                if info: 
                    enhanced += f"\n🌐 Browser: {info.get('title', '?')[:50]}"
                
                if not self.mini_dom:
                    self.mini_dom = self.browser.get_mini_dom() or ""
                
                if self.mini_dom:
                    enhanced += f"\n\n📋 BROWSER-ELEMENTE:\n{self.mini_dom}"
            
            if self.page_text:
                enhanced += f"\n\n📖 SEITEN-TEXT:\n{self.page_text[:2000]}"
            
            if self._file_content:
                enhanced += f"\n\n💻 GELESENE DATEI (✅ ERFOLGREICH):\n{self._file_content[:8000]}"
                enhanced += "\n\n✅ Daten sind sichtbar - du kannst 'done' sagen!"
            elif self._read_ok == False:
                enhanced += "\n\n⚠️ LETZTES read_file FEHLGESCHLAGEN! NICHT 'done' sagen!"
            
            msgs[-1] = {"role": last["role"], "content": enhanced}
            
            t1 = time.time()
            resp, err = call_llm(self.cur, msgs, ss_b64)
            ss_b64 = None
            t2 = time.time()
            
            self.token_count += (sum(len(m.get("content", "")) for m in msgs) + len(resp or "")) // 4
            self.root.after(0, self.update_tokens)
            
            msgs[-1] = last
            
            if err: self.root.after(0, lambda e=err: self.log(f"    ⚠️ {e}\n"))
            self.root.after(0, lambda r=resp, t=t2-t1: self.log(f"    🤖 ({t:.1f}s) {r}\n"))
            
            data = parse_json(resp)
            action = data.get('action', '?')
            
            # ═══════════════════════════════════════════════════════════════
            # DOKUMENT-AKTIONEN (NEU!)
            # ═══════════════════════════════════════════════════════════════
            
            if action == 'create_docx':
                path = os.path.expandvars(data.get('path', '').replace('\\\\', '\\'))
                title = data.get('title', '')
                content = data.get('content', '')
                self.root.after(0, lambda p=path: self.log(f"    📄 Erstelle DOCX: {p}\n"))
                ok, msg = create_docx_file(path, title, content)
                self.root.after(0, lambda ok=ok, m=msg: self.log(f"    {'✅' if ok else '❌'} {m}\n"))
            
            elif action == 'create_xlsx':
                path = os.path.expandvars(data.get('path', '').replace('\\\\', '\\'))
                sheet_data = data.get('data', [])
                self.root.after(0, lambda p=path: self.log(f"    📊 Erstelle XLSX: {p}\n"))
                ok, msg = create_xlsx_file(path, sheet_data)
                self.root.after(0, lambda ok=ok, m=msg: self.log(f"    {'✅' if ok else '❌'} {m}\n"))
            
            elif action == 'create_pptx':
                path = os.path.expandvars(data.get('path', '').replace('\\\\', '\\'))
                slides = data.get('slides', [])
                title = data.get('title', '')
                self.root.after(0, lambda p=path: self.log(f"    📽 Erstelle PPTX: {p}\n"))
                ok, msg = create_pptx_file(path, slides, title)
                self.root.after(0, lambda ok=ok, m=msg: self.log(f"    {'✅' if ok else '❌'} {m}\n"))
            
            # ═══════════════════════════════════════════════════════════════
            # STANDARD-AKTIONEN
            # ═══════════════════════════════════════════════════════════════
            
            elif action == 'mouse_click':
                ok, _ = mouse_click(data.get('x'), data.get('y'), data.get('button', 'left'), 2 if data.get('double') else 1)
                self.root.after(0, lambda: self.log(f"    🖱️ {'✅' if ok else '❌'}\n"))
                time.sleep(0.2)
            
            elif action == 'browser_start':
                url = data.get('url', '')
                self.root.after(0, lambda u=url[:50]: self.log(f"    🌐 Start: {u}\n"))
                try:
                    self.browser.disconnect()
                    self.browser = BrowserHelper()
                    self.mini_dom = ""
                    if IS_WINDOWS:
                        subprocess.run(['taskkill', '/f', '/im', 'chrome.exe'], capture_output=True, 
                                      creationflags=subprocess.CREATE_NO_WINDOW)
                    time.sleep(0.5)
                    cmd_args = [find_chrome_path(), '--remote-debugging-port=9222', '--no-first-run',
                           '--user-data-dir=' + os.path.join(tempfile.gettempdir(), 'chrome-debug')]
                    if url: cmd_args.append(url)
                    subprocess.Popen(cmd_args, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                                    creationflags=subprocess.CREATE_NO_WINDOW if IS_WINDOWS else 0)
                    time.sleep(2)
                    if self.browser.connect():
                        self.root.after(0, lambda: self.log(f"    ✅ Verbunden\n"))
                        time.sleep(1)
                        self.mini_dom = self.browser.get_mini_dom() or ""
                        if self.mini_dom:
                            self.root.after(0, lambda n=len(self.mini_dom.split('\n')): self.log(f"    📋 {n} Elemente im Mini-DOM\n"))
                except Exception as e:
                    self.root.after(0, lambda e=str(e)[:50]: self.log(f"    ❌ {e}\n"))
            
            elif action == 'get_dom':
                self.mini_dom = self.browser.get_mini_dom() or ""
                n = len(self.mini_dom.split('\n')) if self.mini_dom else 0
                self.root.after(0, lambda n=n: self.log(f"    📋 Mini-DOM: {n} Elemente\n"))
            
            elif action == 'playwright_click':
                ok = self.browser.click(data.get('selector'), data.get('text'))
                self.root.after(0, lambda: self.log(f"    🖱️ {'✅' if ok else '❌'}\n"))
                time.sleep(0.3)
                self.mini_dom = self.browser.get_mini_dom() or ""
            
            elif action == 'playwright_type':
                ok = self.browser.type_into(data.get('selector', ''), data.get('text', ''))
                self.root.after(0, lambda: self.log(f"    ⌨️ {'✅' if ok else '❌'}\n"))
            
            elif action == 'playwright_get_text':
                txt = self.browser.get_text(data.get('selector', 'body'))
                if txt:
                    self.page_text = txt
                    self.root.after(0, lambda n=len(txt): self.log(f"    📄 {n} Zeichen\n"))
            
            elif action == 'playwright_navigate':
                if self.browser.navigate(data.get('url', '')):
                    time.sleep(0.5)
                    self.mini_dom = self.browser.get_mini_dom() or ""
                    n = len(self.mini_dom.split('\n')) if self.mini_dom else 0
                    self.root.after(0, lambda n=n: self.log(f"    📋 {n} Elemente\n"))
            
            elif action == 'playwright_scroll':
                self.browser.scroll(data.get('direction', 'down'))
                time.sleep(0.3)
                self.mini_dom = self.browser.get_mini_dom() or ""
            
            elif action == 'key':
                k = data.get('key', '')
                if k:
                    self.root.after(0, lambda k=k: self.log(f"    ⌨️ {k}\n"))
                    press_key(k)
                time.sleep(0.1)
            
            elif action == 'run_commands':
                cmds = data.get('commands', [])
                if cmds:
                    self.root.after(0, lambda n=len(cmds): self.log(f"    💻 {n} Befehle\n"))
                    for c in cmds:
                        self.root.after(0, lambda c=c[:60]: self.log(f"      → {c}...\n"))
                        try:
                            if 'Out-File' in c or ' > ' in c:
                                timeout = 120 if 'powershell' in c.lower() else 60
                                r = subprocess.run(c, shell=True, timeout=timeout, capture_output=True,
                                                  text=True, encoding='utf-8', errors='replace',
                                                  creationflags=subprocess.CREATE_NO_WINDOW if IS_WINDOWS else 0)
                                self.root.after(0, lambda ok=r.returncode==0: self.log(f"      {'✅' if ok else '⚠️'}\n"))
                            else:
                                subprocess.Popen(c, shell=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                                               creationflags=subprocess.CREATE_NO_WINDOW if IS_WINDOWS else 0)
                        except Exception as ex:
                            self.root.after(0, lambda e=str(ex)[:50]: self.log(f"      ❌ {e}\n"))
                time.sleep(0.5)
            
            elif action == 'read_file':
                path = os.path.expandvars(data.get('path', '').replace('\\\\', '\\'))
                self.root.after(0, lambda p=path: self.log(f"    📄 Lese: {p}\n"))
                
                if not os.path.exists(path):
                    for alt in [os.path.join(tempfile.gettempdir(), os.path.basename(path)),
                               os.path.join(tempfile.gettempdir(), 'sysinfo.txt')]:
                        if os.path.exists(alt):
                            path = alt
                            break
                
                if os.path.exists(path):
                    try:
                        content = None
                        for enc in ['utf-8', 'utf-8-sig', 'utf-16', 'cp1252', 'latin-1']:
                            try:
                                with open(path, 'r', encoding=enc) as f:
                                    content = f.read()
                                if content and '\x00' not in content[:100]:
                                    break
                            except: continue
                        
                        if content:
                            self._file_content = content
                            self._read_ok = True
                            def show():
                                self.log("\n" + "═"*70 + "\n")
                                self.log(content)
                                self.log("\n" + "═"*70 + "\n")
                            self.root.after(0, show)
                            self.root.after(0, lambda n=len(content): self.log(f"    ✅ {n} Zeichen\n"))
                        else:
                            self._read_ok = False
                            self.root.after(0, lambda: self.log(f"    ❌ Datei leer\n"))
                    except Exception as ex:
                        self._read_ok = False
                        self.root.after(0, lambda e=str(ex): self.log(f"    ❌ {e}\n"))
                else:
                    self._read_ok = False
                    self.root.after(0, lambda: self.log(f"    ❌ Nicht gefunden\n"))
            
            elif action == 'pywinauto_connect':
                ok = self.pw.connect(data.get('title'), data.get('title_re'))
                self.root.after(0, lambda: self.log(f"    🪟 {'✅' if ok else '❌'}\n"))
            
            elif action == 'pywinauto_type':
                txt = data.get('text', '')
                if txt:
                    self.root.after(0, lambda t=txt[:30]: self.log(f"    ⌨️ '{t}'\n"))
                    try:
                        import pyperclip
                        pyperclip.copy(txt.replace('\\n', '\n'))
                        time.sleep(0.1)
                        pyautogui.hotkey('ctrl', 'v')
                        time.sleep(0.3)
                    except:
                        pyautogui.write(txt.replace('ä', 'ae').replace('ö', 'oe').replace('ü', 'ue'))
                    if data.get('auto_enter'):
                        time.sleep(0.1)
                        press_key('Return')
            
            elif action == 'screenshot':
                b64, _, img = take_screenshot()
                ss_b64 = b64
                self.last_screenshot_img = img
                self.last_screenshot_b64 = b64
                self.screenshot_count += 1
                self.root.after(0, lambda c=self.screenshot_count: self.log(f"    📸 Screenshot #{c} (📋 Zwischenablage)\n"))
            
            elif action == 'wait':
                self.root.after(0, lambda: self.log("    ⏳\n"))
                time.sleep(1)
            
            elif action == 'done':
                self.root.after(0, lambda m=data.get('message', 'Fertig'): self.log(f"\n    ✅ FERTIG: {m}\n"))
                break
            
            else:
                self.root.after(0, lambda a=action: self.log(f"    ⚠️ Unbekannt: {a}\n"))
            
            msgs.append({"role": "assistant", "content": resp})
            msgs.append({"role": "user", "content": "Weiter. Nächster Schritt?"})
            time.sleep(0.05)
        
        self.msgs = msgs
        self.root.after(0, self.finish)

    def finish(self):
        self.run = False
        self.send_btn.config(state='normal')
        self.stop_btn.config(state='disabled')
        self.ent.config(state='normal')
        self.log(f"\n{'='*60}\n📊 Screenshots: {self.screenshot_count}\n")

    def cleanup(self):
        try:
            self.browser.disconnect()
            self.pw.disconnect()
        except: pass
    
    def start(self):
        self.root.protocol("WM_DELETE_WINDOW", lambda: (self.cleanup(), self.root.destroy()))
        self.root.mainloop()


if __name__ == "__main__":
    print("Remote V42 B - KI Remote PC with Playwright + pywinauto + Screenshot")
    print("F1=Stop, 'NEU' am Anfang = neuer Kontext")
    print("NEU: Screenshots automatisch in Zwischenablage!")
    print("     Direkte Dokumenterstellung (create_docx, create_xlsx, create_pptx)")
    
    # Prüfe und installiere Doc-Libraries
    install_doc_libraries()
    
    if not os.path.exists(LLM_CONFIG_DIR):
        os.makedirs(LLM_CONFIG_DIR, exist_ok=True)
    App().start()